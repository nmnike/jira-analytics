"""Smoke-tests for ExportService.

Проверяем, что сгенерированные файлы открываются обратно
соответствующими библиотеками и содержат ожидаемые артефакты
(листы, строки, слайды). Детальный layout не проверяем.
"""

from datetime import datetime
from io import BytesIO

import pytest

from app.models import (
    BacklogItem,
    Category,
    CategoryMapping,
    Employee,
    EmployeeTeam,
    Issue,
    MandatoryWorkType,
    Project,
    Role,
    RoleCapacityRule,
    Worklog,
)
from app.services.categories import CategoryCode
from app.services.export_service import ExportService


@pytest.fixture
def scenario_seed(db_session):
    """Два сотрудника + бэклог + сгенерированный сценарий Q1 2026.

    v3: нужен productive work_type + связанная Category + fallback-правило 100%,
    иначе team_quarter_capacity отдаёт 0 и сценарий не пакует ничего.
    """
    wt = MandatoryWorkType(
        code="productive", label="Продуктив", is_active=True,
    )
    db_session.add(wt)
    db_session.flush()
    db_session.add(
        Category(
            code="cat_productive",
            label="Productive",
            is_system=False,
            work_type_id=wt.id,
        )
    )
    db_session.add(
        RoleCapacityRule(
            year=2026, quarter=1, role=None,
            work_type_id=wt.id, percent_of_norm=100.0,
        )
    )
    db_session.flush()

    # v3 planning per-role: сотрудникам нужна роль из whitelist (analyst/dev/qa),
    # иначе team_role_capacity их пропустит и ёмкость = 0.
    alice = Employee(
        jira_account_id="a1", display_name="Alice", is_active=True, role="dev"
    )
    bob = Employee(
        jira_account_id="b1", display_name="Bob", is_active=True, role="dev"
    )
    db_session.add_all([alice, bob])
    db_session.flush()

    items = [
        BacklogItem(
            title="Redesign login",
            priority=1,
            estimate_hours=100,
            estimate_dev_hours=100,
        ),
        BacklogItem(
            title="Payments v2",
            priority=2,
            estimate_hours=200,
            estimate_dev_hours=200,
        ),
        BacklogItem(
            title="Overflow feature",
            priority=3,
            estimate_hours=5000,
            estimate_dev_hours=5000,
        ),
    ]
    db_session.add_all(items)
    db_session.flush()

    # Manual scenario + allocations (replaces the old greedy generate_scenario).
    from app.models import PlanningScenario, ScenarioAllocation
    db_session.add(Role(
        code="dev", label="Разработчик", color="#1890FF",
        is_active=True, counts_in_planning=True,
    ))
    db_session.add_all([
        EmployeeTeam(employee_id=alice.id, team="DevTeam", is_primary=True),
        EmployeeTeam(employee_id=bob.id, team="DevTeam", is_primary=True),
    ])
    db_session.flush()
    scenario = PlanningScenario(name="Q1 baseline", year=2026, quarter="Q1", status="draft", team="DevTeam")
    db_session.add(scenario)
    db_session.flush()
    # Two items fit (100 + 200 = 300 ≤ 1024), the overflow one doesn't.
    for item, included, planned in [(items[0], True, 100.0),
                                    (items[1], True, 200.0),
                                    (items[2], False, 0.0)]:
        db_session.add(
            ScenarioAllocation(
                scenario_id=scenario.id,
                backlog_item_id=item.id,
                included_flag=included,
                planned_hours=planned,
            )
        )
    db_session.flush()

    class _Result:
        pass
    r = _Result()
    r.scenario_id = scenario.id
    return r


class TestScenarioXlsx:
    def test_workbook_has_four_sheets(self, db_session, scenario_seed):
        from openpyxl import load_workbook

        data = ExportService(db_session).build_scenario_xlsx(
            scenario_seed.scenario_id
        )
        wb = load_workbook(BytesIO(data))
        assert wb.sheetnames == [
            "Сводка", "Включено", "Не вошло", "Справочник",
        ]

    def test_included_titles_present(self, db_session, scenario_seed):
        from openpyxl import load_workbook

        data = ExportService(db_session).build_scenario_xlsx(
            scenario_seed.scenario_id
        )
        wb = load_workbook(BytesIO(data))
        ws_in = wb["Включено"]
        # Header row 2, data starts row 3
        in_titles = {
            ws_in.cell(row=i, column=2).value
            for i in range(3, ws_in.max_row + 1)
        }
        assert "Redesign login" in in_titles
        assert "Payments v2" in in_titles
        assert "Overflow feature" not in in_titles

        ws_out = wb["Не вошло"]
        out_titles = {
            ws_out.cell(row=i, column=2).value
            for i in range(3, ws_out.max_row + 1)
        }
        assert "Overflow feature" in out_titles

    def test_unknown_scenario_raises(self, db_session):
        with pytest.raises(ValueError, match="not found"):
            ExportService(db_session).build_scenario_xlsx("nope")


class TestScenarioPptx:
    def test_presentation_has_four_slides(self, db_session, scenario_seed):
        from pptx import Presentation

        data = ExportService(db_session).build_scenario_pptx(
            scenario_seed.scenario_id
        )
        prs = Presentation(BytesIO(data))
        assert len(prs.slides) == 4

    def test_title_slide_contains_name(self, db_session, scenario_seed):
        from pptx import Presentation

        data = ExportService(db_session).build_scenario_pptx(
            scenario_seed.scenario_id
        )
        prs = Presentation(BytesIO(data))

        title_texts = []
        for shape in prs.slides[0].shapes:
            if shape.has_text_frame:
                title_texts.append(shape.text_frame.text)
        joined = "\n".join(title_texts)
        assert "Q1 baseline" in joined
        assert "Q1" in joined
        assert "2026" in joined

    def test_unknown_scenario_raises(self, db_session):
        with pytest.raises(ValueError, match="not found"):
            ExportService(db_session).build_scenario_pptx("nope")
