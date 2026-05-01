"""Сервис экспортов отчётов в xlsx / pdf / pptx.

Генерит байты готовых файлов. Эндпоинты оборачивают результат в
StreamingResponse с нужным MIME-типом.

Формат экспортов:
- **Аналитика (xlsx)** — многолистовая книга с агрегатами из
  AnalyticsService (сотрудники/проекты/категории/периоды/переключения).
- **Аналитика (pdf)** — сводный отчёт с теми же таблицами, простой layout.
- **Сценарий (xlsx)** — один лист со сводкой сценария и таблицей
  раскладок.
- **Сценарий (pptx)** — несколько слайдов: титул, ключевые метрики,
  таблица включённых задач, пропущенные задачи.

Все модули импортируются лениво внутри методов, чтобы отсутствие одной
из библиотек не ломало импорт остального проекта.
"""

from dataclasses import dataclass
from datetime import datetime
from io import BytesIO
from typing import Optional

from sqlalchemy.orm import Session

from app.models import (
    BacklogItem,
    PlanningScenario,
    ScenarioAllocation,
)
from app.services.analytics_service import AnalyticsService, AggregateRow
from app.services.planning_service import PlanningService


@dataclass
class ScenarioExportRow:
    """Строка экспорта сценария (одна запись раскладки)."""

    title: str
    priority: Optional[int]
    estimate_hours: float
    planned_hours: float
    included: bool


class ExportService:
    """Сервис экспортов аналитики и сценариев планирования."""

    def __init__(self, db: Session):
        self.db = db

    # === Helpers ===

    @staticmethod
    def _fmt_period(start: Optional[datetime], end: Optional[datetime]) -> str:
        if not start and not end:
            return "за всё время"
        s = start.strftime("%Y-%m-%d") if start else "…"
        e = end.strftime("%Y-%m-%d") if end else "…"
        return f"{s} — {e}"

    def _collect_analytics(
        self,
        start: Optional[datetime],
        end: Optional[datetime],
        teams: Optional[list[str]] = None,
        match_employees: bool = True,
        match_issues: bool = True,
    ) -> dict:
        """Собрать все аналитические отчёты за период."""
        analytics = AnalyticsService(self.db)
        kw = dict(
            teams=teams,
            match_employees=match_employees,
            match_issues=match_issues,
        )
        return {
            "by_employee": analytics.hours_by_employee(start, end, **kw),
            "by_project": analytics.hours_by_project(start, end, **kw),
            "by_category": analytics.hours_by_category(start, end, **kw),
            "by_period": analytics.hours_by_period("month", start, end, **kw),
            "switching": analytics.context_switching(start, end, **kw),
        }

    def _load_scenario_rows(
        self, scenario_id: str
    ) -> tuple[PlanningScenario, list[ScenarioExportRow], dict]:
        """Загрузить сценарий и его раскладки, склеив с BacklogItem.

        Возвращает (сценарий, строки_для_экспорта, сводка).
        """
        scenario = self.db.get(PlanningScenario, scenario_id)
        if scenario is None:
            raise ValueError(f"Scenario {scenario_id} not found")

        allocations = (
            self.db.query(ScenarioAllocation, BacklogItem)
            .join(BacklogItem, ScenarioAllocation.backlog_item_id == BacklogItem.id)
            .filter(ScenarioAllocation.scenario_id == scenario_id)
            .all()
        )

        rows = [
            ScenarioExportRow(
                title=item.title,
                priority=item.priority,
                estimate_hours=float(item.estimate_hours or 0.0),
                planned_hours=float(alloc.planned_hours or 0.0),
                included=bool(alloc.included_flag),
            )
            for alloc, item in allocations
        ]
        rows.sort(
            key=lambda r: (
                not r.included,
                r.priority is None,
                r.priority if r.priority is not None else 0,
                r.title,
            )
        )

        total_planned = sum(r.planned_hours for r in rows if r.included)
        included = [r for r in rows if r.included]
        skipped = [r for r in rows if not r.included]
        summary = {
            "total_planned_hours": total_planned,
            "included_count": len(included),
            "skipped_count": len(skipped),
        }
        return scenario, rows, summary

    # === Analytics: Excel ===

    def build_analytics_xlsx(
        self,
        start: Optional[datetime] = None,
        end: Optional[datetime] = None,
        teams: Optional[list[str]] = None,
        match_employees: bool = True,
        match_issues: bool = True,
    ) -> bytes:
        """Собрать многолистовой xlsx со всеми аналитическими отчётами."""
        from openpyxl import Workbook
        from openpyxl.styles import Font, Alignment

        data = self._collect_analytics(
            start, end, teams, match_employees, match_issues
        )

        wb = Workbook()
        bold = Font(bold=True)
        center = Alignment(horizontal="center")

        def _write_sheet(
            ws,
            title: str,
            headers: list[str],
            rows: list[tuple],
        ) -> None:
            ws.title = title
            ws.append(headers)
            for cell in ws[1]:
                cell.font = bold
                cell.alignment = center
            for row in rows:
                ws.append(row)
            for i, _ in enumerate(headers, start=1):
                ws.column_dimensions[
                    ws.cell(row=1, column=i).column_letter
                ].width = 24

        # Первый лист Workbook уже создан
        ws_summary = wb.active
        ws_summary.title = "Сводка"
        ws_summary["A1"] = "Отчёт по аналитике"
        ws_summary["A1"].font = Font(bold=True, size=14)
        ws_summary["A2"] = f"Период: {self._fmt_period(start, end)}"
        ws_summary["A3"] = (
            f"Сформировано: {datetime.now().strftime('%Y-%m-%d %H:%M')}"
        )
        ws_summary.column_dimensions["A"].width = 40

        total_hours = sum(r.total_hours for r in data["by_employee"])
        ws_summary["A5"] = "Всего часов (по сотрудникам):"
        ws_summary["B5"] = round(total_hours, 2)
        ws_summary["A6"] = "Уникальных сотрудников:"
        ws_summary["B6"] = len(data["by_employee"])
        ws_summary["A7"] = "Уникальных проектов:"
        ws_summary["B7"] = len(data["by_project"])
        ws_summary["A8"] = "Категорий:"
        ws_summary["B8"] = len(data["by_category"])

        def _agg_rows(rows: list[AggregateRow]) -> list[tuple]:
            return [
                (r.label, round(r.total_hours, 2), r.worklog_count)
                for r in rows
            ]

        _write_sheet(
            wb.create_sheet("По сотрудникам"),
            "По сотрудникам",
            ["Сотрудник", "Часы", "Worklog"],
            _agg_rows(data["by_employee"]),
        )
        _write_sheet(
            wb.create_sheet("По проектам"),
            "По проектам",
            ["Проект", "Часы", "Worklog"],
            _agg_rows(data["by_project"]),
        )
        _write_sheet(
            wb.create_sheet("По категориям"),
            "По категориям",
            ["Категория", "Часы", "Worklog"],
            _agg_rows(data["by_category"]),
        )
        _write_sheet(
            wb.create_sheet("По месяцам"),
            "По месяцам",
            ["Период", "Часы", "Worklog"],
            _agg_rows(data["by_period"]),
        )
        _write_sheet(
            wb.create_sheet("Переключения"),
            "Переключения",
            [
                "Сотрудник",
                "Worklog",
                "Проектов",
                "Категорий",
                "Переключений",
            ],
            [
                (
                    r.employee_name,
                    r.total_worklogs,
                    r.distinct_projects,
                    r.distinct_categories,
                    r.switches,
                )
                for r in data["switching"]
            ],
        )

        buf = BytesIO()
        wb.save(buf)
        return buf.getvalue()

    # === Analytics: PDF ===

    def build_analytics_pdf(
        self,
        start: Optional[datetime] = None,
        end: Optional[datetime] = None,
        teams: Optional[list[str]] = None,
        match_employees: bool = True,
        match_issues: bool = True,
    ) -> bytes:
        """Собрать PDF-отчёт со сводкой и таблицами."""
        from reportlab.lib import colors
        from reportlab.lib.pagesizes import A4
        from reportlab.lib.styles import getSampleStyleSheet
        from reportlab.platypus import (
            SimpleDocTemplate,
            Paragraph,
            Spacer,
            Table,
            TableStyle,
            PageBreak,
        )

        data = self._collect_analytics(
            start, end, teams, match_employees, match_issues
        )

        buf = BytesIO()
        doc = SimpleDocTemplate(buf, pagesize=A4, title="Jira Analytics")
        styles = getSampleStyleSheet()
        story: list = []

        story.append(Paragraph("Отчёт по аналитике", styles["Title"]))
        story.append(
            Paragraph(
                f"Период: {self._fmt_period(start, end)}",
                styles["Normal"],
            )
        )
        story.append(
            Paragraph(
                f"Сформировано: {datetime.now().strftime('%Y-%m-%d %H:%M')}",
                styles["Normal"],
            )
        )
        story.append(Spacer(1, 12))

        total_hours = sum(r.total_hours for r in data["by_employee"])
        summary_rows = [
            ["Показатель", "Значение"],
            ["Всего часов", f"{total_hours:.2f}"],
            ["Сотрудников", str(len(data["by_employee"]))],
            ["Проектов", str(len(data["by_project"]))],
            ["Категорий", str(len(data["by_category"]))],
        ]
        story.append(self._pdf_table(summary_rows, colors))
        story.append(Spacer(1, 18))

        def _agg_section(title: str, rows: list[AggregateRow]) -> None:
            story.append(Paragraph(title, styles["Heading2"]))
            table_rows = [["Наименование", "Часы", "Worklog"]] + [
                [r.label, f"{r.total_hours:.2f}", str(r.worklog_count)]
                for r in rows
            ]
            if len(table_rows) == 1:
                table_rows.append(["(нет данных)", "", ""])
            story.append(self._pdf_table(table_rows, colors))
            story.append(Spacer(1, 12))

        _agg_section("По сотрудникам", data["by_employee"])
        _agg_section("По проектам", data["by_project"])
        _agg_section("По категориям", data["by_category"])

        story.append(PageBreak())
        _agg_section("По месяцам", data["by_period"])

        story.append(Paragraph("Контекстные переключения", styles["Heading2"]))
        switch_rows = [
            [
                "Сотрудник",
                "Worklog",
                "Проектов",
                "Категорий",
                "Переключений",
            ]
        ] + [
            [
                r.employee_name,
                str(r.total_worklogs),
                str(r.distinct_projects),
                str(r.distinct_categories),
                str(r.switches),
            ]
            for r in data["switching"]
        ]
        if len(switch_rows) == 1:
            switch_rows.append(["(нет данных)", "", "", "", ""])
        story.append(self._pdf_table(switch_rows, colors))

        doc.build(story)
        return buf.getvalue()

    @staticmethod
    def _pdf_table(rows: list[list[str]], colors):
        """Единый стиль PDF-таблицы — шапка серая, сетка чёрная."""
        from reportlab.platypus import Table, TableStyle

        table = Table(rows, repeatRows=1)
        table.setStyle(
            TableStyle(
                [
                    ("BACKGROUND", (0, 0), (-1, 0), colors.lightgrey),
                    ("FONTNAME", (0, 0), (-1, 0), "Helvetica-Bold"),
                    ("GRID", (0, 0), (-1, -1), 0.25, colors.black),
                    ("FONTSIZE", (0, 0), (-1, -1), 9),
                    ("VALIGN", (0, 0), (-1, -1), "MIDDLE"),
                ]
            )
        )
        return table

    # === Capacity: Excel ===

    def export_capacity_xlsx(self, year: int, quarter: int) -> bytes:
        """Excel-выгрузка квартальной ёмкости команды с группировкой по команде."""
        from io import BytesIO
        from openpyxl import Workbook
        from openpyxl.styles import Font, PatternFill, Alignment

        from app.services.capacity_service import CapacityService

        MONTH_LABELS = {
            1: "Янв", 2: "Фев", 3: "Мар", 4: "Апр", 5: "Май", 6: "Июн",
            7: "Июл", 8: "Авг", 9: "Сен", 10: "Окт", 11: "Ноя", 12: "Дек",
        }
        Q = {1: (1, 2, 3), 2: (4, 5, 6), 3: (7, 8, 9), 4: (10, 11, 12)}

        rows = CapacityService(self.db).team_quarter_capacity(year, quarter)

        wb = Workbook()
        ws = wb.active
        ws.title = f"Capacity Q{quarter} {year}"

        bold = Font(bold=True)
        team_fill = PatternFill("solid", fgColor="FFF2E8")
        right = Alignment(horizontal="right")

        months = Q[quarter]
        header: list = ["Сотрудник"]
        for m in months:
            header += [f"{MONTH_LABELS[m]} План", f"{MONTH_LABELS[m]} Факт", f"{MONTH_LABELS[m]} %"]
        header += ["Итого План", "Итого Факт", "Итого %"]
        ws.append(header)
        for c in ws[1]:
            c.font = bold

        groups: dict[str, list] = {}
        for r in rows:
            key = r.team or "__none__"
            groups.setdefault(key, []).append(r)
        ordered_keys = sorted([k for k in groups if k != "__none__"])
        if "__none__" in groups:
            ordered_keys.append("__none__")

        def _pct(plan: float, fact: float) -> str:
            return f"{round(fact / plan * 100)}%" if plan > 0 else "—"

        def _month_value(r, m, attr):
            mc = next((x for x in r.months if x.month == m), None)
            return getattr(mc, attr) if mc else 0.0

        for k in ordered_keys:
            members = groups[k]
            label = "Без команды" if k == "__none__" else k

            team_plan_per_m = [sum(_month_value(r, m, "available_hours") for r in members) for m in months]
            team_fact_per_m = [sum(_month_value(r, m, "fact_hours") for r in members) for m in months]

            team_row: list = [label]
            for plan, fact in zip(team_plan_per_m, team_fact_per_m):
                team_row += [round(plan, 1), round(fact, 1), _pct(plan, fact)]
            total_plan = sum(r.total_available_hours for r in members)
            total_fact = sum(r.total_fact_hours for r in members)
            team_row += [round(total_plan, 1), round(total_fact, 1), _pct(total_plan, total_fact)]
            ws.append(team_row)
            for c in ws[ws.max_row]:
                c.font = bold
                c.fill = team_fill

            for r in members:
                emp_row: list = [r.employee_name]
                for m in months:
                    mc = next((x for x in r.months if x.month == m), None)
                    if mc is None:
                        emp_row += ["", "", ""]
                    else:
                        emp_row += [round(mc.available_hours, 1), round(mc.fact_hours, 1), _pct(mc.available_hours, mc.fact_hours)]
                emp_row += [round(r.total_available_hours, 1), round(r.total_fact_hours, 1), _pct(r.total_available_hours, r.total_fact_hours)]
                ws.append(emp_row)
                for c in ws[ws.max_row][1:]:
                    c.alignment = right

        ws.column_dimensions["A"].width = 28
        for col_letter in "BCDEFGHIJKLM":
            ws.column_dimensions[col_letter].width = 12

        buf = BytesIO()
        wb.save(buf)
        return buf.getvalue()

    # === Analytics Report: Excel ===

    def export_analytics_report_xlsx(
        self,
        report: "AnalyticsReportResponse",
        visible_columns: list[str],
    ) -> bytes:
        """Плоская xlsx-выгрузка иерархического отчёта Аналитики (по задачам)."""
        from io import BytesIO
        from openpyxl import Workbook

        wb = Workbook()
        ws = wb.active
        ws.title = "Аналитика"

        base_headers = [
            "Команда", "Роль", "Сотрудник", "Вид работ", "Категория",
            "Ключ", "Заголовок", "Тип", "Статус", "Часы факт",
        ]
        col_label_map = {
            "plan_hours": "Часы план",
            "pct_plan": "% план",
            "pct_total": "% от итога",
            "worklog_count": "Ворклогов",
            "issue_count": "Задач",
            "employee_count": "Сотрудников",
            "avg_worklog_minutes": "Ср.мин",
            "last_worklog_at": "Последний ворклог",
            "assignee_name": "Исполнитель",
        }
        opt_headers = []
        opt_keys = []
        for col in visible_columns:
            label = col_label_map.get(col)
            if label:
                opt_headers.append(label)
                opt_keys.append(col)

        ws.append(base_headers + opt_headers)

        for team in report.teams:
            for role in team.roles:
                for emp in role.employees:
                    for wt in emp.work_types:
                        for cat in wt.categories:
                            for issue in cat.issues:
                                row: list = [
                                    team.team or "Без команды",
                                    role.role_label,
                                    emp.name,
                                    wt.label,
                                    cat.label,
                                    issue.key,
                                    issue.summary,
                                    issue.issue_type,
                                    issue.status,
                                    issue.totals.fact_hours,
                                ]
                                for k in opt_keys:
                                    if k == "plan_hours":
                                        row.append(issue.totals.plan_hours if issue.totals.plan_hours is not None else "")
                                    elif k == "pct_plan":
                                        row.append(issue.totals.pct_plan if issue.totals.pct_plan is not None else "")
                                    elif k == "pct_total":
                                        row.append(issue.totals.pct_total)
                                    elif k == "worklog_count":
                                        row.append(issue.totals.worklog_count)
                                    elif k == "issue_count":
                                        row.append(issue.totals.issue_count)
                                    elif k == "employee_count":
                                        row.append(issue.totals.employee_count)
                                    elif k == "avg_worklog_minutes":
                                        row.append(issue.totals.avg_worklog_minutes)
                                    elif k == "last_worklog_at":
                                        row.append(issue.last_worklog_at.isoformat() if issue.last_worklog_at else "")
                                    elif k == "assignee_name":
                                        row.append(issue.assignee_name or "")
                                    else:
                                        row.append("")
                                ws.append(row)

        buf = BytesIO()
        wb.save(buf)
        return buf.getvalue()

    # === Scenario: Excel ===

    def build_scenario_xlsx(self, scenario_id: str) -> bytes:
        """Собрать xlsx со сводкой и раскладкой сценария — делегирует в ScenarioXlsxExporter."""
        from app.services.scenario_xlsx_export import ScenarioXlsxExporter

        return ScenarioXlsxExporter(self.db, scenario_id).build()

    # === Scenario: PPTX ===

    def build_scenario_pptx(self, scenario_id: str) -> bytes:
        """Собрать презентацию со сводкой сценария."""
        from pptx import Presentation
        from pptx.util import Inches, Pt

        scenario, rows, summary = self._load_scenario_rows(scenario_id)

        quarter_num = int(scenario.quarter.replace("Q", "")) if scenario.quarter else 0
        total_capacity = 0.0
        if scenario.year and quarter_num:
            total_capacity = PlanningService(self.db)._team_capacity_hours(
                scenario.year, quarter_num
            )
        leftover = max(0.0, total_capacity - summary["total_planned_hours"])

        prs = Presentation()
        blank = prs.slide_layouts[6]

        # --- Slide 1: title ---
        slide = prs.slides.add_slide(blank)
        tb = slide.shapes.add_textbox(Inches(0.5), Inches(1.5), Inches(9), Inches(2))
        frame = tb.text_frame
        p = frame.paragraphs[0]
        p.text = f"Сценарий: {scenario.name}"
        p.font.size = Pt(36)
        p.font.bold = True
        p2 = frame.add_paragraph()
        p2.text = f"{scenario.quarter or ''} {scenario.year or ''}".strip()
        p2.font.size = Pt(20)

        # --- Slide 2: summary metrics ---
        slide = prs.slides.add_slide(blank)
        tb = slide.shapes.add_textbox(Inches(0.5), Inches(0.4), Inches(9), Inches(0.6))
        tb.text_frame.text = "Ключевые метрики"
        tb.text_frame.paragraphs[0].font.size = Pt(28)
        tb.text_frame.paragraphs[0].font.bold = True

        metrics = [
            ("Ёмкость команды, ч", f"{total_capacity:.1f}"),
            ("Запланировано, ч", f"{summary['total_planned_hours']:.1f}"),
            ("Остаток, ч", f"{leftover:.1f}"),
            ("Включено задач", str(summary["included_count"])),
            ("Пропущено задач", str(summary["skipped_count"])),
        ]
        self._pptx_table(
            slide,
            ["Показатель", "Значение"],
            [list(m) for m in metrics],
            top=Inches(1.3),
            width=Inches(6),
        )

        # --- Slide 3: included items ---
        included_rows = [r for r in rows if r.included]
        slide = prs.slides.add_slide(blank)
        tb = slide.shapes.add_textbox(Inches(0.5), Inches(0.4), Inches(9), Inches(0.6))
        tb.text_frame.text = f"Включено в квартал ({len(included_rows)})"
        tb.text_frame.paragraphs[0].font.size = Pt(24)
        tb.text_frame.paragraphs[0].font.bold = True
        self._pptx_table(
            slide,
            ["Задача", "Приоритет", "План, ч"],
            [
                [
                    r.title[:60],
                    str(r.priority) if r.priority is not None else "—",
                    f"{r.planned_hours:.1f}",
                ]
                for r in included_rows
            ]
            or [["(нет задач)", "", ""]],
            top=Inches(1.2),
            width=Inches(9),
        )

        # --- Slide 4: skipped items ---
        skipped_rows = [r for r in rows if not r.included]
        slide = prs.slides.add_slide(blank)
        tb = slide.shapes.add_textbox(Inches(0.5), Inches(0.4), Inches(9), Inches(0.6))
        tb.text_frame.text = f"Не вошло в квартал ({len(skipped_rows)})"
        tb.text_frame.paragraphs[0].font.size = Pt(24)
        tb.text_frame.paragraphs[0].font.bold = True
        self._pptx_table(
            slide,
            ["Задача", "Приоритет", "Оценка, ч"],
            [
                [
                    r.title[:60],
                    str(r.priority) if r.priority is not None else "—",
                    f"{r.estimate_hours:.1f}",
                ]
                for r in skipped_rows
            ]
            or [["(нет задач)", "", ""]],
            top=Inches(1.2),
            width=Inches(9),
        )

        buf = BytesIO()
        prs.save(buf)
        return buf.getvalue()

    @staticmethod
    def _pptx_table(
        slide, headers: list[str], rows: list[list[str]], *, top, width
    ) -> None:
        from pptx.util import Inches, Pt

        n_rows = len(rows) + 1
        n_cols = len(headers)
        max_rows_on_slide = 15
        n_rows = min(n_rows, max_rows_on_slide + 1)

        table_shape = slide.shapes.add_table(
            n_rows,
            n_cols,
            Inches(0.5),
            top,
            width,
            Inches(0.4 * n_rows),
        )
        table = table_shape.table

        for j, h in enumerate(headers):
            cell = table.cell(0, j)
            cell.text = h
            for p in cell.text_frame.paragraphs:
                for run in p.runs:
                    run.font.bold = True
                    run.font.size = Pt(12)

        for i, row in enumerate(rows[: max_rows_on_slide], start=1):
            for j, value in enumerate(row):
                cell = table.cell(i, j)
                cell.text = str(value)
                for p in cell.text_frame.paragraphs:
                    for run in p.runs:
                        run.font.size = Pt(11)
